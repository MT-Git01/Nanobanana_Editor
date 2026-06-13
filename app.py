import streamlit as st
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import os
import io
import json
import tempfile
import fitz
import unicodedata


# =====================================================================
# Page Configuration & Styles
# =====================================================================
st.set_page_config(
    page_title="Nanobanana Infographic Editor",
    page_icon="🍌",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize Session States at the very beginning
if "ocr_blocks" not in st.session_state:
    st.session_state.ocr_blocks = []
if "bg_image" not in st.session_state:
    st.session_state.bg_image = None
if "orig_image" not in st.session_state:
    st.session_state.orig_image = None
if "image_aspect" not in st.session_state:
    st.session_state.image_aspect = 16/9
if "selected_block_id" not in st.session_state:
    st.session_state.selected_block_id = None
if "source_file_name" not in st.session_state:
    st.session_state.source_file_name = ""
if "selected_block_ids" not in st.session_state:
    st.session_state.selected_block_ids = set()

# Read and apply custom CSS styling
if os.path.exists("style.css"):
    with open("style.css", "r", encoding="utf-8") as f:
        css_content = f.read()
        st.markdown(f"<style>{css_content}</style>", unsafe_allow_html=True)

# Custom fonts load (Outfit & Inter)
st.markdown(
    '<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;800&family=Outfit:wght@400;600;800&display=swap" rel="stylesheet">',
    unsafe_allow_html=True
)

# Header Banner
st.markdown(
    """
    <div class="banner">
        <h1 class="gradient-text" style="margin: 0; font-family: 'Outfit', sans-serif; font-size: 2.5rem;">🍌 Nanobanana Editor</h1>
        <p style="margin: 5px 0 0 0; color: #a0aec0; font-size: 1.1rem; font-family: 'Inter', sans-serif;">
            インフォグラフィック画像・PDFを解析し、「文字の選択的消去・背景復元」と「編集可能なスライドテキストオブジェクトへの変換」を同時に行います。
        </p>
    </div>
    """,
    unsafe_allow_html=True
)

# =====================================================================
# Helper Functions
# =====================================================================

def estimate_text_width(text, font_size):
    """Estimate text width in pixels based on font size and CJK/ASCII distribution."""
    width = 0
    for char in text:
        if ord(char) < 128:
            width += font_size * 0.58  # Approximate width of English character
        else:
            width += font_size * 1.0   # Approximate width of CJK character
    return width

def extract_dominant_text_color(image_rgb, x, y, w, h):
    """
    Extract the dominant text color within a bounding box.
    Heuristically filters out the estimated background color and ensures
    appropriate contrast (especially on dark backgrounds).
    """
    img_h, img_w, _ = image_rgb.shape
    x1 = max(0, int(x))
    y1 = max(0, int(y))
    x2 = min(img_w, int(x + w))
    y2 = min(img_h, int(y + h))
    
    crop = image_rgb[y1:y2, x1:x2]
    if crop.size == 0:
        return (255, 255, 255)
    
    # Simple perimeter sampling to estimate the background color
    borders = []
    borders.extend(crop[0, :, :])
    borders.extend(crop[-1, :, :])
    borders.extend(crop[:, 0, :])
    borders.extend(crop[:, -1, :])
    
    borders = np.array(borders)
    median_bg = np.median(borders, axis=0)
    bg_r, bg_g, bg_b = median_bg
    bg_brightness = 0.299 * bg_r + 0.587 * bg_g + 0.114 * bg_b
    
    # Find pixels that differ from the background
    diffs = np.linalg.norm(crop - median_bg, axis=2)
    text_pixels = crop[diffs > 45]
    
    if len(text_pixels) > 0:
        raw_color = np.median(text_pixels, axis=0)
        text_color = tuple(map(int, raw_color))
    else:
        # Fallback to high contrast black or white based on background brightness
        text_color = (0, 0, 0) if bg_brightness > 127 else (255, 255, 255)
        
    # Apply Contrast and Readability Enhancements
    r, g, b = text_color
    if bg_brightness < 120:  # Dark background (expanded threshold)
        # If the extracted text color is too dark, scale up brightness while preserving hue
        max_val = max(r, g, b)
        if max_val < 190:
            if max_val == 0:
                text_color = (255, 255, 255)  # pure white fallback
            else:
                scale = 245 / max_val
                text_color = (min(255, int(r * scale)), min(255, int(g * scale)), min(255, int(b * scale)))
    elif bg_brightness >= 160:  # Light background
        # If the extracted text color is too light, scale down to make it darker
        min_val = min(r, g, b)
        if min_val > 110:
            scale = 70 / min_val
            text_color = (max(0, int(r * scale)), max(0, int(g * scale)), max(0, int(b * scale)))
            
    return text_color


def inpaint_image(image_bgr, blocks, selected_ids):
    """Erase only selected bounding box regions using OpenCV's inpainting."""
    img_h, img_w, _ = image_bgr.shape
    mask = np.zeros((img_h, img_w), dtype=np.uint8)
    
    for block in blocks:
        if block["id"] in selected_ids:
            x, y, w, h = int(block["x"]), int(block["y"]), int(block["width"]), int(block["height"])
            # Add padding/dilation to fully cover anti-aliasing text edges
            padding = 4
            x1 = max(0, x - padding)
            y1 = max(0, y - padding)
            x2 = min(img_w, x + w + padding)
            y2 = min(img_h, y + h + padding)
            cv2.rectangle(mask, (x1, y1), (x2, y2), 255, -1)
            
    if not np.any(mask):
        return image_bgr.copy()
        
    # Inpaint using TELEA method
    inpainted = cv2.inpaint(image_bgr, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)
    return inpainted


def generate_checkerboard(width, height, box_size=12):
    """Generates a gray and white checkerboard pattern image."""
    checkerboard = Image.new("RGBA", (width, height), (255, 255, 255, 255))
    draw = ImageDraw.Draw(checkerboard)
    for y in range(0, height, box_size):
        for x in range(0, width, box_size):
            if ((x // box_size) + (y // box_size)) % 2 == 0:
                draw.rectangle([x, y, x + box_size, y + box_size], fill=(225, 225, 230, 255))
    return checkerboard


def generate_cutout_image(image_pil, blocks, selected_ids):
    """
    Creates an RGBA image where the regions of selected blocks are made transparent,
    and draws outlines around them for user guidance.
    """
    rgba_img = image_pil.convert("RGBA")
    mask = Image.new("L", rgba_img.size, 255)
    mask_draw = ImageDraw.Draw(mask)
    
    for block in blocks:
        if block["id"] in selected_ids:
            x, y, w, h = int(block["x"]), int(block["y"]), int(block["width"]), int(block["height"])
            padding = 4
            x1 = max(0, x - padding)
            y1 = max(0, y - padding)
            x2 = min(rgba_img.width, x + w + padding)
            y2 = min(rgba_img.height, y + h + padding)
            mask_draw.rectangle([x1, y1, x2, y2], fill=0)
            
    rgba_img.putalpha(mask)
    
    # Generate checkerboard background
    checkerboard = generate_checkerboard(rgba_img.width, rgba_img.height)
    
    # Paste transparent image onto checkerboard
    cutout_composited = Image.alpha_composite(checkerboard, rgba_img)
    
    # Draw guides and highlights
    draw = ImageDraw.Draw(cutout_composited)
    for block in blocks:
        x, y, w, h = int(block["x"]), int(block["y"]), int(block["width"]), int(block["height"])
        is_selected = block["id"] in selected_ids
        
        if is_selected:
            # Coral outline for selected cutout regions
            draw.rectangle([x, y, x + w, y + h], outline=(255, 107, 107, 255), width=2)
            # Add a subtle transparent overlay over the cutout area to tint it
            overlay = Image.new("RGBA", cutout_composited.size, (0, 0, 0, 0))
            overlay_draw = ImageDraw.Draw(overlay)
            overlay_draw.rectangle([x, y, x + w, y + h], fill=(255, 107, 107, 30))
            cutout_composited = Image.alpha_composite(cutout_composited, overlay)
            draw = ImageDraw.Draw(cutout_composited)
        else:
            # Yellow outline for potential blocks
            draw.rectangle([x, y, x + w, y + h], outline=(255, 215, 0, 120), width=1)
            
        # Draw ID badge
        label_text = f"#{block['id']}"
        draw.text(
            (x + 4, y + 4),
            label_text,
            fill=(255, 255, 255, 255),
            stroke_fill=(0, 0, 0, 255),
            stroke_width=1
        )
        
    return cutout_composited


def draw_live_preview(bg_image_pil, blocks, selected_ids):
    """Draw edited texts only for selected blocks on the background image."""
    preview_img = bg_image_pil.copy()
    draw = ImageDraw.Draw(preview_img)
    
    for block in blocks:
        if block["id"] not in selected_ids:
            continue
            
        text = block["text_edit"]
        x, y = int(block["x"]), int(block["y"])
        color = block["color"]
        size = block.get("font_size_final", 14)
        
        # Load a suitable system font for rendering preview
        font = None
        font_style = block.get("font_style", "Gothic")
        
        # Candidate font paths (NFC strings)
        candidates = []
        if font_style == "Gothic":
            candidates = [
                "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
                "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc",
                "/System/Library/Fonts/Hiragino Sans GB.ttc",
                "/System/Library/Fonts/Supplemental/Arial.ttf",
                "/Library/Fonts/Arial.ttf"
            ]
        elif font_style == "Mincho":
            candidates = [
                "/System/Library/Fonts/ヒラギノ明朝 ProN.ttc",
                "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
                "/Library/Fonts/Times New Roman.ttf"
            ]
        elif font_style == "Round":
            candidates = [
                "/System/Library/Fonts/ヒラギノ丸ゴ ProN W4.ttc",
                "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
                "/System/Library/Fonts/Hiragino Sans GB.ttc",
                "/System/Library/Fonts/Supplemental/Arial.ttf"
            ]
        else:  # Design / Fallbacks
            candidates = [
                "/System/Library/Fonts/Hiragino Sans GB.ttc",
                "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc",
                "/System/Library/Fonts/Supplemental/Impact.ttf",
                "/Library/Fonts/Impact.ttf"
            ]
            
        # Try loading fonts by normalizing paths to NFD for Mac compat
        for p in candidates:
            try:
                norm_p = unicodedata.normalize('NFD', p)
                font = ImageFont.truetype(norm_p, int(size * 0.95))
                break
            except:
                continue
                
        if font is None:
            font = ImageFont.load_default()
            
        draw.text((x, y), text, fill=color, font=font)
        
    return preview_img


def on_checkbox_change(block_id):
    """Callback function triggered when a checkbox is toggled to update target blocks."""
    cb_key = f"checkbox_block_{block_id}"
    is_checked = st.session_state[cb_key]
    if is_checked:
        st.session_state.selected_block_ids.add(block_id)
    else:
        st.session_state.selected_block_ids.discard(block_id)
        
    # Regenerate selectively-inpainted background image
    orig_img = st.session_state.orig_image
    if orig_img is not None:
        img_np_rgb = np.array(orig_img)
        img_np_bgr = cv2.cvtColor(img_np_rgb, cv2.COLOR_RGB2BGR)
        bg_np_bgr = inpaint_image(img_np_bgr, st.session_state.ocr_blocks, st.session_state.selected_block_ids)
        bg_np_rgb = cv2.cvtColor(bg_np_bgr, cv2.COLOR_BGR2RGB)
        st.session_state.bg_image = Image.fromarray(bg_np_rgb)


def generate_demo_image():
    """Generates a beautiful modern dark-mode infographic demo image."""
    w, h = 1200, 675
    img = Image.new("RGB", (w, h))
    draw = ImageDraw.Draw(img)
    
    # Gradient Background
    for y in range(h):
        r = int(15 + (25 - 15) * (y / h))
        g = int(15 + (35 - 15) * (y / h))
        b = int(25 + (50 - 25) * (y / h))
        draw.line([(0, y), (w, y)], fill=(r, g, b))
        
    # UI Elements & Box Frameworks
    draw.rectangle([50, 40, 1150, 120], fill=(255, 255, 255, 8), outline=(255, 255, 255, 20), width=1)
    
    # Column Panels
    draw.rounded_rectangle([80, 180, 400, 520], radius=16, fill=(26, 28, 45), outline=(52, 58, 90), width=2)
    draw.rounded_rectangle([440, 180, 760, 520], radius=16, fill=(30, 24, 48), outline=(62, 50, 96), width=2)
    draw.rounded_rectangle([800, 180, 1120, 520], radius=16, fill=(22, 32, 45), outline=(44, 66, 92), width=2)
    
    # Graphic Accents
    draw.ellipse([210, 210, 270, 270], fill=(255, 107, 107))
    draw.ellipse([570, 210, 630, 270], fill=(138, 95, 255))
    draw.ellipse([930, 210, 990, 270], fill=(78, 205, 196))
    
    # Title Text
    title_text = "NANOBANANA INFOGRAPHIC EDITOR"
    draw.text((80, 62), title_text, fill=(255, 180, 100))
    
    # Column 1 Text
    draw.text((120, 300), "TEXT INPAINTING", fill=(255, 107, 107))
    draw.text((120, 340), "Erase text and\ninpaint background\nautomatically.", fill=(200, 200, 210))
    
    # Column 2 Text
    draw.text((480, 300), "EDITABLE SLIDE", fill=(138, 95, 255))
    draw.text((480, 340), "Convert text into\nPowerPoint blocks\nfor Google Slides.", fill=(200, 200, 210))
    
    # Column 3 Text
    draw.text((840, 300), "AUTO SCALING", fill=(78, 205, 196))
    draw.text((840, 340), "Adjust sizes and\nprevent layout\noverflow instantly.", fill=(200, 200, 210))
    
    demo_image_path = os.path.join(tempfile.gettempdir(), "demo_infographic.png")
    img.save(demo_image_path)
    
    mock_blocks = [
        {"id": 0, "text": "NANOBANANA INFOGRAPHIC EDITOR", "x": 80, "y": 62, "width": 640, "height": 38, "font_style": "Design", "color": (255, 180, 100)},
        
        {"id": 1, "text": "TEXT INPAINTING", "x": 120, "y": 300, "width": 240, "height": 26, "font_style": "Gothic", "color": (255, 107, 107)},
        {"id": 2, "text": "Erase text and\ninpaint background\nautomatically.", "x": 120, "y": 340, "width": 260, "height": 80, "font_style": "Gothic", "color": (200, 200, 210)},
        
        {"id": 3, "text": "EDITABLE SLIDE", "x": 480, "y": 300, "width": 240, "height": 26, "font_style": "Gothic", "color": (138, 95, 255)},
        {"id": 4, "text": "Convert text into\nPowerPoint blocks\nfor Google Slides.", "x": 480, "y": 340, "width": 260, "height": 80, "font_style": "Gothic", "color": (200, 200, 210)},
        
        {"id": 5, "text": "AUTO SCALING", "x": 840, "y": 300, "width": 240, "height": 26, "font_style": "Gothic", "color": (78, 205, 196)},
        {"id": 6, "text": "Adjust sizes and\nprevent layout\noverflow instantly.", "x": 840, "y": 340, "width": 260, "height": 80, "font_style": "Gothic", "color": (200, 200, 210)}
    ]
    return demo_image_path, mock_blocks


# =====================================================================
# Google Cloud Vision OCR Logic
# =====================================================================
def run_cloud_vision_ocr(image_bytes):
    """Connect to Google Cloud Vision API and parse paragraph text boxes."""
    try:
        from google.cloud import vision
        client = vision.ImageAnnotatorClient()
        image = vision.Image(content=image_bytes)
        response = client.document_text_detection(image=image)
        
        if response.error.message:
            raise Exception(response.error.message)
            
        blocks_data = []
        block_id = 0
        
        for page in response.full_text_annotation.pages:
            for block in page.blocks:
                for paragraph in block.paragraphs:
                    # Stitch text characters
                    para_text = ""
                    for word in paragraph.words:
                        word_text = "".join([symbol.text for symbol in word.symbols])
                        para_text += word_text
                        
                        # Handle breaks (spaces, newlines)
                        last_symbol = word.symbols[-1]
                        if last_symbol.property.detected_break:
                            break_type = last_symbol.property.detected_break.type_
                            if break_type in (1, 2):  # SPACE, SURE_SPACE
                                para_text += " "
                            elif break_type in (3, 5):  # EOL, LINE_BREAK
                                para_text += "\n"
                                
                    # Coordinate calculation
                    vertices = paragraph.bounding_box.vertices
                    xs = [v.x for v in vertices if v.x is not None]
                    ys = [v.y for v in vertices if v.y is not None]
                    
                    if xs and ys:
                        x_min, x_max = min(xs), max(xs)
                        y_min, y_max = min(ys), max(ys)
                        width = x_max - x_min
                        height = y_max - y_min
                        
                        # Only keep text boxes with actual printable characters
                        clean_text = para_text.strip()
                        if clean_text:
                            blocks_data.append({
                                "id": block_id,
                                "text": clean_text,
                                "x": x_min,
                                "y": y_min,
                                "width": width,
                                "height": height,
                                "font_style": "Gothic",  # default guessed font
                            })
                            block_id += 1
                            
        return blocks_data, None
    except Exception as e:
        return None, str(e)


# =====================================================================
# Sidebar Configuration
# =====================================================================
st.sidebar.markdown(
    """
    <div style='text-align: center; padding-bottom: 15px;'>
        <h2 style='margin: 0; font-family: "Outfit", sans-serif; color: #ff6b6b;'>⚙️ Settings</h2>
    </div>
    """,
    unsafe_allow_html=True
)

demo_mode = st.sidebar.toggle("デモモードで試す (Mock OCR)", value=True)

# Google Credentials Config
gcp_cred_path = None
if not demo_mode:
    st.sidebar.markdown("---")
    st.sidebar.subheader("🔑 Google Cloud Credentials")
    gcp_file = st.sidebar.file_uploader("Credentials JSON File", type=["json"])
    
    if gcp_file is not None:
        try:
            creds = json.load(gcp_file)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".json") as tmp:
                tmp.write(json.dumps(creds).encode('utf-8'))
                gcp_cred_path = tmp.name
            os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = gcp_cred_path
            st.sidebar.success("✅ Google Credentials Load Successful!")
        except Exception as e:
            st.sidebar.error(f"❌ JSONファイル読み込みエラー: {e}")
    elif "GOOGLE_APPLICATION_CREDENTIALS" in os.environ:
        st.sidebar.info("💡 システム環境変数から資格情報を検出しました。")
    else:
        st.sidebar.warning("⚠️ Google Cloud Vision API 資格情報が未設定です。")

st.sidebar.markdown("---")
st.sidebar.subheader("📐 出力スライド比率")
slide_aspect = st.sidebar.selectbox(
    "アスペクト比を選択",
    ["16:9 (Widescreen)", "4:3 (Standard)"]
)

if "16:9" in slide_aspect:
    slide_width_inch = 13.333
    slide_height_inch = 7.5
else:
    slide_width_inch = 10.0
    slide_height_inch = 7.5


# =====================================================================
# Application Main Workflow
# =====================================================================

# Handle Source Selection
if demo_mode:
    # Set up Demo Environment
    if st.session_state.source_file_name != "demo" or not st.session_state.ocr_blocks:
        demo_img_path, mock_blocks = generate_demo_image()
        orig_img = Image.open(demo_img_path)
        img_np_bgr = cv2.imread(demo_img_path)
        img_np_rgb = cv2.cvtColor(img_np_bgr, cv2.COLOR_BGR2RGB)
        
        # Populate session state
        st.session_state.orig_image = orig_img
        st.session_state.image_aspect = orig_img.width / orig_img.height
        st.session_state.source_file_name = "demo"
        st.session_state.selected_block_ids = {0, 1, 3, 5}
        
        # Add scaling parameters to mock blocks
        blocks = []
        for b in mock_blocks:
            lines = b["text"].split("\n")
            num_lines = max(1, len([l for l in lines if l.strip()]))
            line_height = b["height"] / num_lines
            font_size_init = max(10, int(line_height * 0.85))
            
            blocks.append({
                "id": b["id"],
                "text": b["text"],
                "text_edit": b["text"],
                "x": b["x"],
                "y": b["y"],
                "width": b["width"],
                "height": b["height"],
                "font_style": b["font_style"],
                "font_size_init": font_size_init,
                "font_size_modifier": 0,
                "color": b["color"]
            })
        st.session_state.ocr_blocks = blocks
        
        # Inpaint background based on default selection
        bg_np_bgr = inpaint_image(img_np_bgr, blocks, st.session_state.selected_block_ids)
        bg_np_rgb = cv2.cvtColor(bg_np_bgr, cv2.COLOR_BGR2RGB)
        st.session_state.bg_image = Image.fromarray(bg_np_rgb)

else:
    # Live Upload UI Panel
    st.markdown("### 📥 インフォグラフィック ファイルアップロード")
    uploaded_file = st.file_uploader(
        "画像またはPDFファイルを選択してください",
        type=["png", "jpg", "jpeg", "pdf"]
    )
    
    if uploaded_file is not None:
        # Avoid re-running OCR on same file
        if st.session_state.source_file_name != uploaded_file.name:
            file_bytes = uploaded_file.read()
            st.session_state.source_file_name = uploaded_file.name
            
            # Check PDF
            if uploaded_file.name.endswith(".pdf"):
                try:
                    doc = fitz.open(stream=file_bytes, filetype="pdf")
                    if len(doc) == 0:
                        st.error("PDF内にページが見つかりませんでした。")
                        st.stop()
                    # Render the first page at high quality (2.0x zoom/144 DPI)
                    page = doc.load_page(0)
                    zoom = 2.0
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")
                    orig_img = Image.open(io.BytesIO(img_data)).convert("RGB")
                except Exception as e:
                    st.error(f"PDF処理エラー: {e}")
                    st.info("💡 PDF解析エラーが発生しました。ファイルが破損していないか、画像形式（PNG/JPG）でお試しください。")
                    st.stop()
            else:
                orig_img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
                
            # Convert PIL image to BGR numpy array for OpenCV
            img_np_rgb = np.array(orig_img)
            img_np_bgr = cv2.cvtColor(img_np_rgb, cv2.COLOR_RGB2BGR)
            
            # Auto-resize huge images to fit memory and render faster
            h, w = img_np_bgr.shape[:2]
            scale_ratio = 1.0
            if max(h, w) > 1920:
                scale_ratio = 1920.0 / max(h, w)
                new_w = int(w * scale_ratio)
                new_h = int(h * scale_ratio)
                img_np_bgr = cv2.resize(img_np_bgr, (new_w, new_h), interpolation=cv2.INTER_AREA)
                img_np_rgb = cv2.cvtColor(img_np_bgr, cv2.COLOR_BGR2RGB)
                orig_img = Image.fromarray(img_np_rgb)
                st.info(f"💡 高解像度画像を最適なパフォーマンスサイズ（{new_w}x{new_h}）に自動縮小しました。")
                
            # Connect Vision API OCR (Supporting Workload Identity and Local ADC natively)
            with st.spinner("🔍 Google Cloud Vision API を使って文字と座標を検出中..."):
                # Convert processed numpy image back to bytes for Google API
                _, buffer = cv2.imencode(".png", img_np_bgr)
                raw_blocks, ocr_error = run_cloud_vision_ocr(buffer.tobytes())
                
            if ocr_error is not None:
                st.error("⛔ Google Cloud Vision API 認証エラーまたは接続エラーが発生しました。")
                st.markdown(
                    f"""
                    **エラー内容:** `{ocr_error}`
                    
                    **💡 解決方法:**
                    1. **Credentials JSON ファイルのアップロード (推奨):**
                       左側のサイドバーを開き、Google Cloud Vision 認証用 JSON ファイルをアップロードしてください。
                    2. **Workload Identity / 環境認証情報 (ADC) の利用:**
                       Cloud Run、GKE、またはローカルPC環境で、Workload Identity Federation 連携設定、または `gcloud auth application-default login` が適切に行われているかご確認ください。環境変数 `GOOGLE_APPLICATION_CREDENTIALS` またはメタデータサーバーから自動で認証が行われます。
                    """
                )
                st.stop()
                
            if raw_blocks:
                # Add metadata & colors
                blocks = []
                for b in raw_blocks:
                    text_color = extract_dominant_text_color(img_np_rgb, b["x"], b["y"], b["width"], b["height"])
                    
                    # Initial font size approximation based on text height and lines
                    lines = b["text"].split("\n")
                    num_lines = max(1, len([l for l in lines if l.strip()]))
                    line_height = b["height"] / num_lines
                    font_size_init = max(10, int(line_height * 0.85))
                    
                    blocks.append({
                        "id": b["id"],
                        "text": b["text"],
                        "text_edit": b["text"],
                        "x": b["x"],
                        "y": b["y"],
                        "width": b["width"],
                        "height": b["height"],
                        "font_style": "Gothic",
                        "font_size_init": font_size_init,
                        "font_size_modifier": 0,
                        "color": text_color
                    })
                st.session_state.ocr_blocks = blocks
                st.session_state.selected_block_ids = set()
                st.session_state.selected_block_id = None
                st.session_state.orig_image = orig_img
                st.session_state.bg_image = orig_img
                st.session_state.image_aspect = orig_img.width / orig_img.height
                st.toast("🚀 画像解析が完了しました！編集したいブロックを選択してください。")
            else:
                st.warning("⚠️ 画像から文字を抽出できませんでした。文字が鮮明であるか、ファイルが壊れていないか確認してください。")
                st.stop()


# =====================================================================
# Main Workspace Presentation Layout
# =====================================================================

if st.session_state.orig_image is not None and st.session_state.ocr_blocks:
    col_canvas, col_editor = st.columns([1, 1])
    
    # Left Column: Cutout Preview & Checkbox Selection List
    with col_canvas:
        st.markdown(
            """
            <div class="glass-panel">
                <h3 style='margin-top:0; color: #ff8e53; font-family: "Outfit", sans-serif;'>🖼️ 切り抜きプレビュー & 編集対象選択</h3>
                <p style='color: #a0aec0; font-size:0.9rem;'>
                    スライド上で編集可能にするテキストブロックにチェックを入れてください。<br>
                    選択されたテキストは画像から透過で切り抜かれ（穴が開き）、右側のエディタで再編集・配置できます。
                </p>
            </div>
            """, 
            unsafe_allow_html=True
        )
        
        # Render transparent cutout image
        cutout_img = generate_cutout_image(
            st.session_state.orig_image,
            st.session_state.ocr_blocks,
            st.session_state.selected_block_ids
        )
        st.image(cutout_img, use_column_width=True)
        
        # Checkbox list in scrollable container
        st.markdown("##### 🔍 検出されたテキストブロック一覧")
        with st.container(height=350):
            for block in st.session_state.ocr_blocks:
                cb_key = f"checkbox_block_{block['id']}"
                is_checked = block["id"] in st.session_state.selected_block_ids
                
                # Checkbox to add/remove block from active list
                st.checkbox(
                    f"#{block['id']}: {block['text']}",
                    value=is_checked,
                    key=cb_key,
                    on_change=on_checkbox_change,
                    args=(block['id'],)
                )
        
    # Right Column: Interactive Editor Dashboard & Selective Live Preview
    with col_editor:
        st.markdown(
            """
            <div class="glass-panel">
                <h3 style='margin-top:0; color: #4facfe; font-family: "Outfit", sans-serif;'>✍️ 編集仕上がりプレビュー (Live Preview)</h3>
                <p style='color: #a0aec0; font-size:0.9rem;'>選択されたテキストが消去され、入力された新しいテキストが重ね合わされたリアルタイムのプレビューです。</p>
            </div>
            """, 
            unsafe_allow_html=True
        )
        
        # Selective live preview rendering
        live_preview_img = draw_live_preview(
            st.session_state.bg_image,
            st.session_state.ocr_blocks,
            st.session_state.selected_block_ids
        )
        st.image(live_preview_img, use_column_width=True)
        
        st.markdown("### 📝 選択テキストの編集")
        selected_list = sorted(list(st.session_state.selected_block_ids))
        if not selected_list:
            st.info("💡 左側の一覧から編集したいテキストを選択してください。")
        else:
            # Dropdown to select which of the selected blocks to customize
            active_id = st.selectbox(
                "編集・カスタマイズするブロックを選択",
                options=selected_list,
                format_func=lambda bid: f"ブロック #{bid}: {st.session_state.ocr_blocks[bid]['text_edit'][:30]}...",
                key="active_block_selector"
            )
            
            block = st.session_state.ocr_blocks[active_id]
            st.markdown(f"#### ⚙️ ブロック #{active_id} 編集設定")
            
            with st.container():
                st.markdown('<div class="block-card">', unsafe_allow_html=True)
                
                # Text content
                new_text_val = st.text_area(
                    "テキスト内容",
                    value=block["text_edit"],
                    key=f"text_edit_val_{active_id}"
                )
                block["text_edit"] = new_text_val
                
                # Font style
                font_options = [
                    "Gothic系 (Arial / MS Gothic)",
                    "Mincho系 (Times New Roman / MS Mincho)",
                    "Round系 (Arial Rounded / HG Maru Gothic)",
                    "Design系 (Impact / Trebuchet MS)"
                ]
                current_style = block["font_style"]
                style_index = 0
                if current_style == "Mincho":
                    style_index = 1
                elif current_style == "Round":
                    style_index = 2
                elif current_style == "Design":
                    style_index = 3
                    
                selected_style = st.selectbox(
                    "フォントファミリー代替案",
                    options=font_options,
                    index=style_index,
                    key=f"font_style_select_{active_id}"
                )
                
                if "Gothic" in selected_style:
                    block["font_style"] = "Gothic"
                    block["font_name"] = "Arial"
                elif "Mincho" in selected_style:
                    block["font_style"] = "Mincho"
                    block["font_name"] = "Times New Roman"
                elif "Round" in selected_style:
                    block["font_style"] = "Round"
                    block["font_name"] = "Arial Rounded MT Bold"
                elif "Design" in selected_style:
                    block["font_style"] = "Design"
                    block["font_name"] = "Impact"
                
                # Custom override
                custom_override = st.checkbox(
                    "カスタムフォント名を入力する (パワーユーザー向け)", 
                    value=("font_custom_name" in block),
                    key=f"custom_font_check_{active_id}"
                )
                if custom_override:
                    custom_font_name = st.text_input(
                        "フォント名 (例: Noto Sans JP)", 
                        value=block.get("font_custom_name", block.get("font_name", "Arial")),
                        key=f"custom_font_val_{active_id}"
                    )
                    block["font_custom_name"] = custom_font_name
                    block["font_name"] = custom_font_name
                elif "font_custom_name" in block:
                    del block["font_custom_name"]
                
                # Auto scale indicator
                w_box = block["width"]
                w_text = estimate_text_width(block["text_edit"], block["font_size_init"])
                scale = min(1.0, w_box / w_text) if w_text > 0 else 1.0
                auto_font_size = int(block["font_size_init"] * scale)
                
                st.markdown(
                    f"""
                    <div style='font-size: 0.85rem; color: #a0aec0; margin-bottom: 8px;'>
                        📐 Bounding Box 幅: <b>{w_box}px</b> | 計算されたテキスト幅: <b>{int(w_text)}px</b><br>
                        🎯 自動縮小倍率: <b>{scale:.2f}x</b> | 自動計算フォントサイズ: <b>{auto_font_size}pt</b>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
                
                # Size modifier slider
                size_modifier = st.slider(
                    "フォントサイズ微調整 (pt)",
                    min_value=-30,
                    max_value=30,
                    value=block["font_size_modifier"],
                    step=1,
                    key=f"font_size_slider_{active_id}"
                )
                block["font_size_modifier"] = size_modifier
                
                final_font_size = max(5, auto_font_size + size_modifier)
                block["font_size_final"] = final_font_size
                
                st.markdown(f"<b>最終適用フォントサイズ:</b> <span style='color: #4facfe; font-size:1.1rem;'>{final_font_size}pt</span>", unsafe_allow_html=True)
                
                # Color picker
                hex_init = '#{:02x}{:02x}{:02x}'.format(*block["color"])
                selected_hex = st.color_picker(
                    "テキストカラー",
                    value=hex_init,
                    key=f"color_picker_{active_id}"
                )
                hex_clean = selected_hex.lstrip('#')
                block["color"] = tuple(int(hex_clean[i:i+2], 16) for i in (0, 2, 4))
                
                st.markdown('</div>', unsafe_allow_html=True)

        # Presentation building & output slide download
        st.markdown("---")
        st.subheader("📦 PowerPoint (.pptx) 出力生成")
        
        # Summary metrics
        st.markdown(
            f"""
            - 検出テキストブロック数: **{len(st.session_state.ocr_blocks)}** 個
            - 編集可能に設定したブロック数: **{len(st.session_state.selected_block_ids)}** 個
            - 出力スライド比率: **{slide_aspect}**
            - 背景画像: **選択された文字のみ消去した背景（OpenCV Inpaint）** を全画面適用
            """
        )
        
        # Build PPTX Presentation dynamically
        try:
            prs = Presentation()
            prs.slide_width = Inches(slide_width_inch)
            prs.slide_height = Inches(slide_height_inch)
            
            # Add slide (blank layout)
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Save background image to a temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_bg:
                st.session_state.bg_image.save(tmp_bg.name)
                bg_path = tmp_bg.name
                
            # Fullbleed background image placement
            slide.shapes.add_picture(bg_path, 0, 0, prs.slide_width, prs.slide_height)
            
            # Overlay editable text frames
            orig_w, orig_h = st.session_state.orig_image.width, st.session_state.orig_image.height
            
            for b in st.session_state.ocr_blocks:
                if b["id"] not in st.session_state.selected_block_ids:
                    continue
                # Calculate slide proportional coordinates
                left = Inches((b["x"] / orig_w) * slide_width_inch)
                top = Inches((b["y"] / orig_h) * slide_height_inch)
                width = Inches((b["width"] / orig_w) * slide_width_inch)
                height = Inches((b["height"] / orig_h) * slide_height_inch)
                
                tx_box = slide.shapes.add_textbox(left, top, width, height)
                tf = tx_box.text_frame
                tf.word_wrap = True
                # Zero out default PPTX margins to match original positioning exactly
                tf.margin_left = Inches(0)
                tf.margin_right = Inches(0)
                tf.margin_top = Inches(0)
                tf.margin_bottom = Inches(0)
                
                # Split text by newline and add as separate paragraphs to maintain exact structure
                lines = b["text_edit"].split("\n")
                
                p = tf.paragraphs[0]
                p.text = lines[0]
                p.font.name = b.get("font_name", "Arial")
                p.font.size = Pt(b.get("font_size_final", 14))
                p.font.color.rgb = RGBColor(*b["color"])
                
                for line in lines[1:]:
                    p_line = tf.add_paragraph()
                    p_line.text = line
                    p_line.font.name = b.get("font_name", "Arial")
                    p_line.font.size = Pt(b.get("font_size_final", 14))
                    p_line.font.color.rgb = RGBColor(*b["color"])
                    
            # Clean up temporary background file
            try:
                os.unlink(bg_path)
            except:
                pass
            
            # Write presentation bytes
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()
        except Exception as e:
            st.error(f"PowerPointスライド生成エラー: {e}")
            pptx_bytes = b""
            
        # Build PNG / JPEG image bytes dynamically
        try:
            preview_img = draw_live_preview(st.session_state.bg_image, st.session_state.ocr_blocks, st.session_state.selected_block_ids)
            
            png_buffer = io.BytesIO()
            preview_img.save(png_buffer, format="PNG")
            png_bytes = png_buffer.getvalue()
            
            jpeg_buffer = io.BytesIO()
            rgb_preview = preview_img.convert("RGB")
            rgb_preview.save(jpeg_buffer, format="JPEG", quality=95)
            jpeg_bytes = jpeg_buffer.getvalue()
        except Exception as e:
            st.error(f"画像書き出しエラー: {e}")
            png_bytes = b""
            jpeg_bytes = b""

        # Direct save to local Downloads folder fallback
        downloads_dir = os.path.expanduser("~/Downloads")
        if os.path.exists(downloads_dir):
            if st.button("💾 ローカルPCの「ダウンロード」フォルダに直接保存する", key="direct_save_btn", use_container_width=True):
                try:
                    pptx_path = os.path.join(downloads_dir, "nanobanana_editable_slide.pptx")
                    png_path = os.path.join(downloads_dir, "nanobanana_edited_image.png")
                    jpeg_path = os.path.join(downloads_dir, "nanobanana_edited_image.jpg")
                    
                    with open(pptx_path, "wb") as f:
                        f.write(pptx_bytes)
                    with open(png_path, "wb") as f:
                        f.write(png_bytes)
                    with open(jpeg_path, "wb") as f:
                        f.write(jpeg_bytes)
                        
                    st.success(f"🎉 ダウンロードフォルダに直接保存しました！\n- {pptx_path}\n- {png_path}\n- {jpeg_path}")
                except Exception as e:
                    st.error(f"直接保存に失敗しました: {e}")

        # Direct download buttons
        st.download_button(
            label="📥 編集可能な.pptxをダウンロード",
            data=pptx_bytes,
            file_name="nanobanana_editable_slide.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            key="download_pptx_btn"
        )
        
        col_png, col_jpg = st.columns(2)
        with col_png:
            st.download_button(
                label="🖼️ 編集後画像 (PNG) をダウンロード",
                data=png_bytes,
                file_name="nanobanana_edited_image.png",
                mime="image/png",
                use_container_width=True,
                key="download_png_btn"
            )
        with col_jpg:
            st.download_button(
                label="🖼️ 編集後画像 (JPEG) をダウンロード",
                data=jpeg_bytes,
                file_name="nanobanana_edited_image.jpg",
                mime="image/jpeg",
                use_container_width=True,
                key="download_jpg_btn"
            )
else:
    # Initial instruction screen if no files uploaded/demo initialized
    st.info("💡 アプリケーションを開始するには、デモモードを有効にするか、画像をアップロードしてください。")

# Clean up GCP Credential temporary file if it was created
if gcp_cred_path and os.path.exists(gcp_cred_path):
    try:
        os.unlink(gcp_cred_path)
    except:
        pass
